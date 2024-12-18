import os
import json
import boto3
import logging
import time
import io
import base64
import concurrent.futures
from .presentation_decomposer import PresentationProcessor, CustomEncoder
from pptx import Presentation

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

class PresentationDecomposer:
    def __init__(self):
        self.processor = PresentationProcessor()
        self.s3_client = self.create_s3_client()
        self.bucket_name = os.environ.get('S3_BUCKET_NAME')

    def create_s3_client(self):
        return boto3.client('s3',
            aws_access_key_id=os.environ.get('AWS_ACCESS_KEY_ID'),
            aws_secret_access_key=os.environ.get('AWS_SECRET_ACCESS_KEY'),
            region_name=os.environ.get('AWS_REGION', 'us-east-1')
        )

    def upload_file(self, image_bytes, s3_key):
        if image_bytes is None:
            logging.warning(f"Skipping upload for {s3_key}: No image data")
            return None
        try:
            self.s3_client.put_object(Bucket=self.bucket_name, Key=s3_key, Body=image_bytes, ContentType='image/png')
            return f"https://{self.bucket_name}.s3.amazonaws.com/{s3_key}"
        except Exception as e:
            logging.error(f"Error uploading file {s3_key}: {str(e)}")
            return None

    def process_presentation(self, pptx_file, slide_index):
        logging.info(f"Starting to process presentation: {pptx_file}")

        total_slides = len(Presentation(pptx_file).slides)
        if slide_index >= total_slides:
            raise ValueError(f"Slide index {slide_index} is out of bounds for the total number of slides in the presentation: {total_slides}")
        
        slide_data = self.processor.process_presentation(pptx_file, slide_index)
        print(slide_data)
        # Ensure slide_data is a dictionary
        if not isinstance(slide_data, dict):
            logging.error(f"Unexpected data type for slide_data: {type(slide_data)}")
            

        # Process and upload shape images
        shapes = {}
        for shape_name, image_bytes in slide_data.get("images", {}).items():
            if image_bytes is not None:
                s3_key = f"processed/{os.path.basename(pptx_file)}/slide_{slide_index}/{shape_name}.png"
                public_url = self.upload_file(image_bytes, s3_key)
                if public_url:
                    shapes[shape_name] = public_url
        # Upload background image
        background_url = None
        if slide_data.get("background"):
            background_s3_key = f"processed/{os.path.basename(pptx_file)}/slide_{slide_index}/background.png"
            background_url = self.upload_file(slide_data["background"], background_s3_key)

        slide_structure = {
            "index": slide_index,
            "shapes": shapes,
            "structure": slide_data.get("structure", {}),
            "thumbnail": self.upload_file(slide_data.get("thumbnail"), f"processed/{os.path.basename(pptx_file)}/slide_{slide_index}/thumbnail.png"),
            "background": background_url,
            "frame_size": slide_data.get("frame_size")  # Add this line
        }

        return json.dumps(slide_structure, cls=CustomEncoder)

    def process_single_slide(self, pptx_file, slide_index):
        start_time = time.time()
        
        logging.info(f"Starting to process slide {slide_index} from presentation: {pptx_file}")
        
        process_start = time.time()
        document = self.processor.open_presentation(pptx_file)
        pptx = Presentation(pptx_file)
        slide_data = self.processor.process_slide(document, pptx, slide_index)
        process_end = time.time()
        logging.info(f"Slide processing took {process_end - process_start:.2f} seconds")
        
        upload_start = time.time()
        with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
            upload_tasks = []

            for shape_name, image_bytes in slide_data['images'].items():
                if image_bytes is not None:
                    s3_key = f"processed/{os.path.basename(pptx_file)}/slide_{slide_index}/{shape_name}.png"
                    task = executor.submit(self.upload_file, image_bytes, s3_key)
                    upload_tasks.append((shape_name, task))

            for shape_name, task in upload_tasks:
                public_url = task.result()
                if public_url:
                    slide_data['images'][shape_name] = public_url

        upload_end = time.time()
        logging.info(f"Uploading files took {upload_end - upload_start:.2f} seconds")

        # Convert result to JSON
        json_result = json.dumps(slide_data, cls=CustomEncoder, indent=2)

        end_time = time.time()
        total_time = end_time - start_time
        logging.info(f"Total processing time: {total_time:.2f} seconds")

        return json_result

import sys
import os
import uno
from unohelper import systemPathToFileUrl
from com.sun.star.beans import PropertyValue
from com.sun.star.awt import Size, Point
import json
from pptx import Presentation
from pptx.util import Emu
from PIL import Image
import uuid
import tempfile
import subprocess
from pdf2image import convert_from_path
import urllib.parse
import io
import logging
from concurrent.futures import ThreadPoolExecutor
import time
import enum
import math

# Add these imports from structure.py and decompose_images.py
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PresentationProcessor:
    def __init__(self):
        self.desktop = self.connect_to_libreoffice()
        self.dpi = 96  # Standard screen DPI
        self.frame_width = None
        self.frame_height = None

    @staticmethod
    def connect_to_libreoffice():
        local_context = uno.getComponentContext()
        resolver = local_context.ServiceManager.createInstanceWithContext(
            "com.sun.star.bridge.UnoUrlResolver", local_context)
        context = resolver.resolve(
            "uno:socket,host=localhost,port=2002;urp;StarOffice.ComponentContext")
        desktop = context.ServiceManager.createInstanceWithContext("com.sun.star.frame.Desktop", context)
        return desktop

    def open_presentation(self, file_path):
        url = uno.systemPathToFileUrl(os.path.abspath(file_path))
        properties = []
        document = self.desktop.loadComponentFromURL(url, "_blank", 0, tuple(properties))
        if not document:
            raise Exception(f"Failed to open document: {file_path}")
        return document

    @staticmethod
    def rename_shapes(pptx_file):
        logger.info(f"Renaming shapes in {pptx_file}")
        import time
        t1 = time.time()
        prs = Presentation(pptx_file)
        shape_counter = 1
        for slide in prs.slides:
            for shape in slide.shapes:
                shape.name = f"Shape_{shape_counter}"
                shape_counter += 1
        prs.save(pptx_file)
        t2 = time.time()
        logger.info(f"Renaming shapes took {t2 - t1:.2f} seconds")
        return pptx_file

    def process_presentation(self, pptx_file):
        document = self.open_presentation(pptx_file)
        pptx = Presentation(pptx_file)
        
        # Get the frame size
        self.frame_width = self.emu_to_pixels(pptx.slide_width)
        self.frame_height = self.emu_to_pixels(pptx.slide_height)
        
        slides = document.getDrawPages()
        total_slides = slides.getCount()

        shape_counter = 1  # Global shape counter

        for slide_index in range(total_slides):
            slide_data = self.process_slide(document, pptx, slide_index, shape_counter)
            shape_counter = slide_data['next_shape_counter']  # Update the global counter
            
            # Add frame size to slide_data
            slide_data['frame_size'] = {
                'width': self.frame_width,
                'height': self.frame_height
            }
            
            yield slide_data

        document.close(True)

    def process_slide(self, document, pptx, slide_index, shape_counter):
        slide = document.getDrawPages().getByIndex(slide_index)
        pptx_slide = pptx.slides[slide_index]


        # Extract components sequentially
        structure_result = self.extract_slide_structure(slide, pptx_slide, shape_counter)
        images_result = self.extract_slide_images(document, slide, pptx_slide, slide_index, shape_counter)
        thumbnail = self.extract_slide_image(document, slide_index)

        # Create a BytesIO object for the background
        background_buffer = io.BytesIO()
        background_buffer = self.extract_background(document, slide, pptx_slide, background_buffer)

        slide_data = {
            "structure": structure_result['structure'],
            "images": images_result['images'],
            "thumbnail": thumbnail,
            "background": background_buffer,
            "next_shape_counter": max(structure_result['next_shape_counter'], images_result['next_shape_counter'])
        }

        return slide_data

    def extract_slide_structure(self, slide, pptx_slide, shape_counter):
        structure = []

        lo_shapes = list(slide)
        pptx_shapes = list(pptx_slide.shapes)

        for lo_shape, pptx_shape in zip(lo_shapes, pptx_shapes):
            shape_name = f"Shape_{shape_counter}"
            shape_counter += 1
            shape_data = self.get_shape_info(lo_shape, pptx_shape, shape_name)
            structure.append(shape_data)

        return {"structure": structure, "next_shape_counter": shape_counter}

    def get_shape_info(self, lo_shape, pptx_shape, shape_name):
        # Get position and size from LibreOffice shape
        position = lo_shape.Position
        size = lo_shape.Size

        height = self.mm100_to_pixels(size.Height)
        width = self.mm100_to_pixels(size.Width)

        _type = str(lo_shape.ShapeType)

        if _type == "com.sun.star.drawing.ConnectorShape" and width == 0:
            width = 2
        
        if _type == "com.sun.star.drawing.ConnectorShape" and height == 0:
            height = 2

        shape_info = {
            "name": shape_name,
            "type": str(lo_shape.ShapeType),
            "width": {"value": width, "unit": "px"},
            "height": {"value": height, "unit": "px"},
            "position_x": {"value": self.mm100_to_pixels(position.X), "unit": "px"},
            "position_y": {"value": self.mm100_to_pixels(position.Y), "unit": "px"},
            "rotation": {"value": int(getattr(lo_shape, "RotateAngle", 0)), "unit": "0.01degree"},
            "z_order": getattr(lo_shape, "ZOrder", None),
            "has_text": False,  # Initialize as False
            "vertical_alignment": self.get_vertical_alignment(lo_shape),
            "text_auto_grow_height": self.get_boolean_attribute(lo_shape, "TextAutoGrowHeight"),
            "text_auto_grow_width": self.get_boolean_attribute(lo_shape, "TextAutoGrowWidth"),
            "text_word_wrap": self.get_boolean_attribute(lo_shape, "TextWordWrap"),
            "text_fit_to_size": self.get_boolean_attribute(lo_shape, "TextFitToSize"),
            "padding": {
                "left": {"value": self.mm100_to_pixels(getattr(lo_shape, "TextLeftDistance", 0)), "unit": "px"},
                "right": {"value": self.mm100_to_pixels(getattr(lo_shape, "TextRightDistance", 0)), "unit": "px"},
                "top": {"value": self.mm100_to_pixels(getattr(lo_shape, "TextUpperDistance", 0)), "unit": "px"},
                "bottom": {"value": self.mm100_to_pixels(getattr(lo_shape, "TextLowerDistance", 0)), "unit": "px"},
            }
        }

        if hasattr(lo_shape, "Text"):
            text_frame = self.extract_text_frame(lo_shape, pptx_shape, shape_info["vertical_alignment"])
            shape_info["text_frame"] = text_frame

            # Determine if the shape actually contains any non-empty text
            # has_text is True only if there's at least one non-empty run in any paragraph
            shape_info["has_text"] = any(
                run["text"].strip()
                for paragraph in text_frame["paragraphs"]
                for run in paragraph["runs"]
            )

            # Comment explaining the ethos of has_text
            """
            has_text Ethos:
            The 'has_text' property indicates whether a shape contains any actual, 
            non-empty text content. It's not enough for a shape to have a text frame; 
            it must have at least one non-empty run of text within its paragraphs.
            This approach ensures that we don't treat shapes with empty text frames 
            or only whitespace as text-containing shapes. It allows for more accurate 
            representation and processing of shapes that genuinely contain textual content.
            """

        shape_info["fill"] = self.extract_fill_properties(lo_shape)
        shape_info["line"] = self.extract_line_properties(lo_shape)

        return shape_info

    def extract_text_frame(self, lo_shape, pptx_shape, vertical_alignment):
        text_frame = {
            "paragraphs": [],
            "vertical_alignment": vertical_alignment
        }
        text = lo_shape.Text
        text_enum = text.createEnumeration()

        # Get the overall text frame position and size
        text_frame_pos_x = self.emu_to_pixels(pptx_shape.left)
        text_frame_pos_y = self.emu_to_pixels(pptx_shape.top)
        text_frame_width = self.emu_to_pixels(pptx_shape.width)
        text_frame_height = self.emu_to_pixels(pptx_shape.height)

        # Initialize variables for paragraph positioning
        current_y = 0
        line_spacing = 1.2  # Default line spacing factor


        bullet_indexes = {i: 0 for i in range(10)} 

        while text_enum.hasMoreElements():
            paragraph = text_enum.nextElement()
            para_info = self.extract_paragraph_info(paragraph, lo_shape, pptx_shape, text_frame_pos_x, text_frame_pos_y + current_y, text_frame_width, line_spacing, bullet_indexes)
            text_frame["paragraphs"].append(para_info)
            current_y += para_info["height"]["value"]

        return text_frame

    def extract_paragraph_info(self, paragraph, lo_shape, pptx_shape, base_x, base_y, width, line_spacing, bullet_indexes):
        para_text = paragraph.getString().strip()
        
        para_info = {
            "text": para_text,
            "alignment": self.get_alignment_description(getattr(paragraph, "ParaAdjust", None)),
            "level": getattr(paragraph, "NumberingLevel", None),
            "bullet_info": self.extract_bullet_info(paragraph, bullet_indexes),
            "runs": [],
            "position_x": {"value": base_x, "unit": "px"},
            "position_y": {"value": base_y, "unit": "px"},
            "width": {"value": width, "unit": "px"},
            "margin": {
                "left": {"value": self.mm100_to_pixels(getattr(paragraph, "ParaLeftMargin", 0)), "unit": "px"},
                "right": {"value": self.mm100_to_pixels(getattr(paragraph, "ParaRightMargin", 0)), "unit": "px"},
                "top": {"value": self.mm100_to_pixels(getattr(paragraph, "ParaTopMargin", 0)), "unit": "px"},
                "bottom": {"value": self.mm100_to_pixels(getattr(paragraph, "ParaBottomMargin", 0)), "unit": "px"},
            },
            "padding": {
                "left": {"value": self.mm100_to_pixels(getattr(lo_shape, "TextLeftDistance", 0)), "unit": "px"},
                "right": {"value": self.mm100_to_pixels(getattr(lo_shape, "TextRightDistance", 0)), "unit": "px"},
                "top": {"value": self.mm100_to_pixels(getattr(lo_shape, "TextUpperDistance", 0)), "unit": "px"},
                "bottom": {"value": self.mm100_to_pixels(getattr(lo_shape, "TextLowerDistance", 0)), "unit": "px"},
            },
            "autofit": self.get_autofit_property(lo_shape),
            "background_color": self.get_rgba_color(getattr(paragraph, "ParaBackColor", None)),
        }

        # Process character runs
        portions = paragraph.createEnumeration()
        total_height = 0
        max_font_size = 0
        while portions.hasMoreElements():
            portion = portions.nextElement()
            run_info = self.extract_run_info(portion, base_x, base_y + total_height)
            para_info["runs"].append(run_info)
            total_height += run_info["height"]["value"]
            max_font_size = max(max_font_size, run_info["font_size"])

        # Calculate line height and paragraph height
        line_height = self.pt_to_pixels(max_font_size) * line_spacing
        para_height = max(total_height, line_height)
        
        para_info["line_height"] = {"value": line_height, "unit": "px"}
        para_info["height"] = {"value": para_height, "unit": "px"}

        return para_info

    def extract_bullet_info(self, paragraph, bullet_indexes):
        try:
            numbering_level = getattr(paragraph, "NumberingLevel", None)
            numbering_rules = getattr(paragraph, "NumberingRules", None)

            bullet_info = {
                "visible": False,
                "type": 0,
                "char": None,
                "image_url": None,
                "bullet_size_percent": None,
                "bullet_color": None,
                "color": None,
                "font_name": None,
                "font_size": {"value": 0, "unit": "px"},
                "bold": False,
                "italic": False,
                "underline": False,
                "orientation": 0,
                "kerning": False,
                "distance": {"value": 0, "unit": "px"},
                "start_value": 0,
                "current_value": 0,
                "indent": {"value": 0, "unit": "px"},
                "prefix": "",
                "suffix": "",
                "level": None,
                "style": "",
                "type_description": "UNKNOWN"
            }

            if numbering_level is not None and numbering_level >= 0:
                bullet_info["visible"] = True
                bullet_info["level"] = numbering_level
                current_index = bullet_indexes[numbering_level]
                bullet_info["current_value"] = current_index
                bullet_indexes[numbering_level] += 1
                for i in range(numbering_level + 1, 10):
                    bullet_indexes[i] = 0
                rule = numbering_rules.getByIndex(numbering_level)
                for property_value in rule:
                    
                    if property_value.Name == "StartWith":
                        bullet_info["start_value"] = property_value.Value
                    
                    if property_value.Name == "LeftMargin":
                        bullet_info["indent"] = {"value": self.mm100_to_pixels(property_value.Value), "unit": "px"}
                    
                    if property_value.Name == "BulletColor":
                        bullet_info["bullet_color"] = self.get_rgba_color(property_value.Value)

                    if property_value.Name == "BulletSizePercent":
                        bullet_info["bullet_size_percent"] = property_value.Value

                    if property_value.Name == "NumberingType":
                        bullet_info["type"] = property_value.Value
                    
                    if property_value.Name == "Prefix":
                        bullet_info["prefix"] = property_value.Value
                    
                    if property_value.Name == "Suffix":
                        bullet_info["suffix"] = property_value.Value
                    
                    if property_value.Name == "BulletChar":
                        print(f"Bullet char: {property_value.Value}")
                        bullet_info["char"] = property_value.Value
                    
                    name = property_value.Name
                    value = property_value.Value
                    state = property_value.State

                    # print(f"Name: {name}, Value: {value}")

                    if name == "BulletFont":
                        bullet_font_name = value.Name
                        bullet_height = value.Height
                        bullet_width = value.Width
                        bullet_weight = value.Weight
                        bullet_underline = value.Underline
                        # bullet_color = value.Color
                        bullet_orientation = value.Orientation
                        bullet_kerning = value.Kerning

                        if value.Name:
                            bullet_info["font_name"] = value.Name
                        if value.Height:
                            bullet_info["font_size"] = {"value": self.pt_to_pixels(value.Height), "unit": "px"}
                        if value.Weight:
                            bullet_info["bold"] = value.Weight > 100
                        if value.Underline:
                            bullet_info["underline"] = value.Underline != 0
                        if value.Orientation:
                            bullet_info["orientation"] = value.Orientation
                        if value.Kerning:
                            bullet_info["kerning"] = value.Kerning
                    print("")

                    if name == "GraphicObject" and value:
                        print("Bullet graphic found!")
                    

            numbering_type_map = {
                0: "NONE",
                1: "OUTLINE",
                2: "BITMAP",
                3: "CHAR_SPECIAL",
                4: "NUMBER",
                5: "BITMAP_SPECIAL",
                6: "CHAR",
                7: "CONTINUE",
                8: "CHAR_UPPER_LETTER",
                9: "CHAR_LOWER_LETTER",
                10: "CHAR_UPPER_ROMAN",
                11: "CHAR_LOWER_ROMAN",
                12: "PAGE_DESCRIPTOR",
                13: "CHAPTER_NUMBER",
                14: "BITMAP_URL"
            }

            bullet_info["type_description"] = numbering_type_map.get(bullet_info["type"], "UNKNOWN")

            return bullet_info

        except Exception as e:
            print(f"Error extracting bullet info: {str(e)}")
            return {
                "visible": False,
                "type": 0,
                "char": None,
                "image_url": None,
                "bullet_size_percent": None,
                "bullet_color": None,
                "color": None,
                "font_name": None,
                "font_size": {"value": 0, "unit": "px"},
                "bold": False,
                "italic": False,
                "underline": False,
                "orientation": 0,
                "kerning": False,
                "distance": {"value": 0, "unit": "px"},
                "start_value": 0,
                "current_value": 0,
                "indent": {"value": 0, "unit": "px"},
                "prefix": "",
                "suffix": "",
                "level": None,
                "style": "",
                "type_description": "UNKNOWN"
            }



    def extract_run_info(self, portion, base_x, base_y):
        font_size = getattr(portion, "CharHeight", None)
        run_text = portion.getString()
        run_width = self.pt_to_pixels(font_size * len(run_text) * 0.6) if font_size else 0
        run_height = self.pt_to_pixels(font_size) if font_size else 0

        return {
            "text": run_text,
            "font_name": getattr(portion, "CharFontName", None),
            "font_size": font_size,
            "bold": getattr(portion, "CharWeight", None) > 100,
            "italic": getattr(portion, "CharPosture", None) == 2,
            "underline": getattr(portion, "CharUnderline", None) != 0,
            "color": self.get_rgba_color(getattr(portion, "CharColor", None)),
            "highlight_color": self.get_rgba_color(getattr(portion, "CharBackColor", None)),
            "position_x": {"value": base_x, "unit": "px"},
            "position_y": {"value": base_y, "unit": "px"},
            "width": {"value": run_width, "unit": "px"},
            "height": {"value": run_height, "unit": "px"},
        }

    def extract_fill_properties(self, shape):
        fill =  {
            "type": getattr(shape, "FillStyle", None),
            "color": self.get_rgba_color(getattr(shape, "FillColor", None)),
            "transparency": {"value": getattr(shape, "FillTransparence", 0), "unit": "%"}
        }

        if fill['color'] and fill['color']['red'] == 114 and fill['color']['green'] == 159 and fill['color']['blue'] == 207:
            fill['color']['alpha'] = 0

        return fill

    def extract_line_properties(self, shape):
        return {
            "color": self.get_rgba_color(getattr(shape, "LineColor", None)),
            "width": {"value": self.mm100_to_pixels(getattr(shape, "LineWidth", 0)), "unit": "px"},
            "style": getattr(shape, "LineStyle", None),
            "transparency": {"value": getattr(shape, "LineTransparence", 0), "unit": "%"}
        }

    def extract_slide_images(self, document, slide, pptx_slide, slide_index, shape_counter):
        images = {}
        
        lo_shapes = list(slide)
        pptx_shapes = list(pptx_slide.shapes)

        for lo_shape, pptx_shape in zip(lo_shapes, pptx_shapes):
            shape_name = f"Shape_{shape_counter}"
            shape_counter += 1
            
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                temp_filename = temp_file.name

            try:
                self.export_shape_as_png(lo_shape, temp_filename, scale_factor=2)
                with open(temp_filename, 'rb') as f:
                    images[shape_name] = io.BytesIO(f.read())
            except Exception as e:
                logger.warning(f"Error exporting shape {shape_name}: {str(e)}")
                print(e)
                print("calling in extract slide images")
            finally:
                if os.path.exists(temp_filename):
                    os.remove(temp_filename)

        return {"images": images, "next_shape_counter": shape_counter}

    def extract_background(self, document, slide, pptx_slide, output_buffer):
        # Store original shapes
        original_shapes = list(slide)
        
        # Remove all shapes from the slide
        while slide.getCount() > 0:
            slide.remove(slide.getByIndex(0))
        
        try:
            # Take a snapshot of the slide (now only containing the background)
            controller = document.getCurrentController()
            controller.setCurrentPage(slide)
            
            export_properties = (
                PropertyValue("FilterName", 0, "impress_png_Export", 0),
                PropertyValue("Overwrite", 0, True, 0),
                PropertyValue("PixelWidth", 0, 2048, 0),
                PropertyValue("PixelHeight", 0, 1536, 0),
            )
            
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                temp_filename = temp_file.name
            
            document.storeToURL(uno.systemPathToFileUrl(temp_filename), export_properties)
            
            with Image.open(temp_filename) as img:
                img.save(output_buffer, format='PNG')
                output_buffer.seek(0)
            
            os.unlink(temp_filename)
        
        finally:
            # Restore original shapes
            for shape in original_shapes:
                slide.add(shape)
        
        return output_buffer

    def get_lo_background_color(self, slide):
        # Try to get background color from LibreOffice slide
        try:
            fill_color = slide.Background.FillColor
            return (fill_color & 0xFF, (fill_color >> 8) & 0xFF, (fill_color >> 16) & 0xFF)
        except AttributeError:
            return (255, 255, 255)  # Default to white if color can't be determined

    def get_pptx_background_color(self, slide):
        # Get background color from pptx slide
        if slide.background.fill.type == MSO_FILL.SOLID:
            color = slide.background.fill.fore_color.rgb
            return (color.red, color.green, color.blue)
        else:
            return (255, 255, 255)  # Default to white if not a solid color

    def export_shape_as_png(self, shape, export_path, scale_factor=2):
        # Rotate the shape to 0 degrees
        try:
            original_rotation = shape.getPropertyValue("RotateAngle")
        except Exception as e:
            logger.warning(f"Error getting rotation: {str(e)}")
            original_rotation = 0
        try:
            shape.setPropertyValue("RotateAngle", 0)
        except Exception as e:
            logger.warning(f"Error setting rotation to 0: {str(e)}")

        export_properties = (
            PropertyValue("FilterName", 0, "draw_png_Export", 0),
            PropertyValue("Overwrite", 0, True, 0),
            PropertyValue("SelectionOnly", 0, True, 0),
        )

        output_url = systemPathToFileUrl(export_path)

        controller = self.desktop.getCurrentFrame().getController()
        selection_supplier = controller.queryInterface(
            uno.getTypeByName('com.sun.star.view.XSelectionSupplier'))
        selection_supplier.select(shape)

        document = self.desktop.getCurrentComponent()
        if not document:
            raise Exception("No document is currently open")
        document.storeToURL(output_url, export_properties)

        # Restore the original rotation
        try:
            shape.setPropertyValue("RotateAngle", original_rotation)
        except Exception as e:
            logger.warning(f"Error restoring rotation: {str(e)}")

        # Resize the exported image using Pillow
        with Image.open(export_path) as img:
            width, height = img.size
            new_size = (int(width * scale_factor), int(height * scale_factor))
            resized_img = img.resize(new_size, Image.LANCZOS)
            resized_img.save(export_path)

    def get_alignment_description(self, alignment_value):
        alignment_map = {
            0: "left",
            1: "right",
            2: "justify",
            3: "center",
            4: "distributed"
        }
        return alignment_map.get(alignment_value, "unknown")

    def get_vertical_alignment(self, shape):
        try:
            vert_align = shape.TextVerticalAdjust
            
            # Map the enum values to our desired strings
            alignment_map = {
                'TOP': "top",
                'CENTER': "center",
                'BOTTOM': "bottom"
            }
            
            # Extract the enum value (e.g., 'TOP') from the string representation
            enum_value = str(vert_align).split("'")[1]
            
            return alignment_map.get(enum_value, "unknown")
        except AttributeError:
            return "unknown"

    def get_rgba_color(self, color):
        if color is None:
            return None
        if isinstance(color, int):
            alpha = 255 - ((color >> 24) & 255)  # Extract alpha from the highest 8 bits
            return {
                "red": (color >> 16) & 255,
                "green": (color >> 8) & 255,
                "blue": color & 255,
                "alpha": alpha
            }
        elif hasattr(color, 'Red') and hasattr(color, 'Green') and hasattr(color, 'Blue'):
            # For LibreOffice color objects, we need to check if there's a Transparency attribute
            transparency = getattr(color, 'Transparency', 0)
            alpha = 255 - int(transparency * 2.55)  # Convert percentage to 0-255 range
            return {
                "red": color.Red,
                "green": color.Green,
                "blue": color.Blue,
                "alpha": alpha
            }
        return None

    def mm100_to_pixels(self, mm100):
        mm = mm100 / 100
        inches = mm / 25.4
        return round(inches * self.dpi)

    def pt_to_pixels(self, pt):
        return round(pt * self.dpi / 72)

    def emu_to_pixels(self, emu):
        return Emu(emu).inches * self.dpi

    def extract_slide_image(self, document, slide_index):
        controller = document.getCurrentController()
        
        # Get the draw pages directly from the document
        pages = document.getDrawPages()
        page = pages.getByIndex(slide_index)
        
        # Set the view to the current slide
        controller.setCurrentPage(page)
        
        # Export properties for PNG with higher resolution
        export_properties = (
            PropertyValue("FilterName", 0, "impress_png_Export", 0),
            PropertyValue("Overwrite", 0, True, 0),
            PropertyValue("PixelWidth", 0, 2048, 0),  # Increased width
            PropertyValue("PixelHeight", 0, 1536, 0),  # Increased height
        )
        
        # Create a temporary file to store the slide image
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
            temp_filename = temp_file.name

        # Export the slide as PNG
        document.storeToURL(uno.systemPathToFileUrl(temp_filename), export_properties)
        
        # Open the exported image using Pillow
        with Image.open(temp_filename) as img:
            # Save the image to a BytesIO object
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format='PNG')
            img_byte_arr.seek(0)

        # Remove the temporary file
        os.unlink(temp_filename)

        return img_byte_arr

    def get_line_spacing_info(self, line_spacing):
        if line_spacing is None:
            return None
        height_mm100 = getattr(line_spacing, "Height", None)
        height_px = self.mm100_to_pixels(height_mm100) if height_mm100 is not None else None
        return {
            "Mode": getattr(line_spacing, "Mode", None),
            "Height": {"value": height_px, "unit": "px"}
        }

    def get_boolean_attribute(self, obj, attr_name):
        value = getattr(obj, attr_name, None)
        if isinstance(value, bool):
            return value
        elif isinstance(value, str):
            return value.lower() == 'true'
        elif isinstance(value, int):
            return value != 0
        elif isinstance(value, uno.Enum):
            # Handle Enum types
            enum_value = str(value).split("'")[1]
            return enum_value.lower() == 'true'
        else:
            return False

    def get_autofit_property(self, lo_shape):
        try:
            fit_to_size = getattr(lo_shape, "TextFitToSize", None)
            auto_grow_height = getattr(lo_shape, "TextAutoGrowHeight", None)
            auto_grow_width = getattr(lo_shape, "TextAutoGrowWidth", None)
            # Handle Enum for TextFitToSize
            if isinstance(fit_to_size, uno.Enum):
                fit_to_size = str(fit_to_size).split("'")[1]
            if fit_to_size == 'NONE' and auto_grow_height == False and auto_grow_width == False:
                return "Do Not Autofit"
            elif fit_to_size == 'NONE' and (auto_grow_height == True or auto_grow_width == True):
                return "Shrink Text on Overflow"
            elif fit_to_size in ['ALLLINES', 'FIRSTLINE']:
                return "Resize Shape to Fit Text"
            else:
                return "Unknown"
        except AttributeError as e:
            return "Unknown"

class CustomEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, io.BytesIO):
            return obj.getvalue().decode('latin1')
        elif isinstance(obj, enum.Enum):
            return obj.name
        elif hasattr(obj, 'value'):
            return obj.value
        elif hasattr(obj, '__dict__'):
            return obj.__dict__
        else:
            return str(obj)



















